import collections
import collections.abc
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '각 팀별 데이터'

#클래스 객체 생성
chart_data = ChartData()

# chart_data.categories 속성
# 원형 차트의 범주를 정의
chart_data.categories = ['총무팀', '영업팀', '개발팀', '운영팀', '마케팅팀']

# 모든 지역의 데이터로 chart_data 객체를 구성
# chart_data.add_series('팀별 비율', (0.135, 0.324, 0.180, 0.235, 0.126))
chart_data.add_series('팀별 비율', (13, 32, 18, 23, 12))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
#프레젠테이션 슬라이드에 차트 추가. 차트의 유형 등을 설정
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart
#데이터 레전드(범례 혹은 라벨) 설정
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
# data_labels.number_format = '0%'
data_labels.number_format = '0\%'
# data_labels.show_percentage = True
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

prs.save('charts_test.pptx')