import collections
import collections.abc
from pptx import Presentation

prs = Presentation()
two_content_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(two_content_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = '2개의 컨텐츠 슬라이드 추가하기'

tf = body_shape.text_frame
tf.text = '첫번째 라인입니다.'

p = tf.add_paragraph()
p.text = '두번째 라인입니다.'
p.level = 1

p = tf.add_paragraph()
p.text = '세번째 라인입니다.'
p.level = 2

prs.save('placeholders_test.pptx')

