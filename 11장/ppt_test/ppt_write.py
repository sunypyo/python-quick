import collections
import collections.abc
from pptx import Presentation

# 빈 템플릿으로 프리젠테이션 객체를 생성한다
prs = Presentation()
# 프리젠테이션에 슬라이드를 추가한다
# 슬라이드 레이아웃[0]은 제목(프레젠테이션 제목 슬라이드) 형식이다
slide = prs.slides.add_slide(prs.slide_layouts[0])
# 슬라이드 쉐이프의 타이틀을 입력한다
slide.shapes.title.text = "안녕, 파이썬 :D"
# 슬라이드 개체틀에 텍스트를 입력한다
slide.placeholders[1].text = "파이썬은 정말 멋져요~"
# 프리젠테이션 파일을 저장한다
prs.save('ppt_write_result.pptx')

# template.pptx 파일을 템플릿으로 하여 프리젠테이션 객체를 생성한다
prs = Presentation('template.pptx')
# 프리젠테이션의 첫번째 슬라이드를 얻어온다
first_slide = prs.slides[0]
# 슬라이드의 첫번째 쉐이프의 텍스트 프레임 첫 단락에 텍스트를 추가한다
first_slide.shapes[0].text_frame.paragraphs[0].text = "첫번째 슬라이드"
# 슬라이드의 두번째 쉐이프의 텍스트 프레임 첫 단락에 텍스트를 추가한다
first_slide.shapes[1].text_frame.paragraphs[0].text = " 파이썬 ppt 다루기 "
# 슬라이드의 두번째 쉐이프의 텍스트 프레임에
# 새로운 단락을 추가하고 텍스트를 추가한다
first_slide.shapes[1].text_frame.add_paragraph().text = " ppt 파일을 생성 "

# 프리젠테이션에 슬라이드를 추가한다
# 슬라이드 레이아웃[1]은 제목 및 내용 형식이다
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
# 단락에 기록할 데이터들을 리스트로 선언한다
datas = [
    ' python-pptx 라이브러리 사용',
    ' Presentation 생성',
    ' Slide 추가',
    ' 단락 추가와 텍스트 추가'
]

# 슬라이드의 첫번째 쉐이프의 텍스트 프레임 첫 단락에 텍스트를 추가한다
slide2.shapes[0].text_frame.paragraphs[0].text = "두번째 슬라이드"

# 슬라이드의 두번째 쉐이프의 텍스트 프레임을 얻어온다
text_frame = slide2.shapes[1].text_frame
# 텍스트 프레임의 첫번째 단락을 꺼낸다
p = text_frame.paragraphs[0]
# 단락의 텍스트에 datas 리스트의 첫번째 문자열 값을 대입한다
p.text = datas[0]

# 위의 코드를 반복문으로 처리한다
# 단락의 텍스트에 datas 리스트의 두번째 문자열 부터 끝까지 대입한다
for data in datas[1:]:
    text_frame.add_paragraph().text = data
# 프리젠테이션 파일을 저장한다
prs.save('ppt_write_result2.pptx')