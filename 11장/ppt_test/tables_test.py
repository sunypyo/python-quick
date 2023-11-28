import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
shapes.title.text = '사원 정보'

table = shapes.add_table(4, 3, Inches(2.0), Inches(2.0), Inches(6.0), Inches(1.2)).table

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(2.0)

table.cell(0, 0).text = '사원 번호'
table.cell(0, 1).text = '사원 이름'
table.cell(0, 2).text = '사원 아이디'

students = {
    1: ["표선영", 115],
    2: ["홍길동", 119],
    3: ["전우치", 101]
}

for i in range(len(students)):
    table.cell(i+1, 0).text = str(i+1)
    table.cell(i+1, 1).text = str(students[i+1][0])
    table.cell(i+1, 2).text = str(students[i+1][1])

prs.save('table_test.pptx')
