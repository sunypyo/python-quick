import collections
import collections.abc
from pptx import Presentation

ppt_file = 'ppt_read_test.pptx'
prs = Presentation(ppt_file)

# 읽어온 텍스트를 저장할 파일 객체 생성
wfile = open('ppt_read_text.txt', mode='w', encoding='utf-8')

# 읽어온 텍스트를 저장할 리스트
texts = [ ]
# 전체 Presentation 객체로부터 모든 슬라이드를 얻은 후, 슬라이드 하나씩 꺼내기
for slide in prs.slides:
    # 하나의 슬라이드 안의 여러 shapes 들을 얻은 후, shape  하나씩 꺼내기
    for shape in slide.shapes:
        # 텍스트가 없다면 반복문을 계속
        if not shape.has_text_frame:
            continue
        # 텍스트가 있다면 텍스트의 paragraphs 들을 얻은 후, paragraph 하나씩 꺼내기
        for paragraph in shape.text_frame.paragraphs:
            # 읽어온 단락의 텍스트 추출
            text = paragraph.text
            # 읽어온 단락의 텍스트들을 화면에 출력
            print(text)
            # 읽어온 단락을 리스트에 추가
            texts.append(text)
            # 읽어온 텍스트들을 파일에 저장
            wfile.write(text + '\n')

wfile.close()

# 구분선과 읽어온 텍스트가 저장된 리스트 출력
print("=" * 30)
print("읽어온 텍스트를 리스트로 보기 : ", texts)

# 구분선과 슬라이드 레이아웃 정보 출력
print("=" * 30)
slide1, slide2 = prs.slides[0], prs.slides[1]
print("슬라이드 레이아웃 정보: \n", slide1.slide_layout.name, ",", slide2.slide_layout.name)
print("슬라이드 shape 정보")

# 슬라이드 쉐이프(모양) 정보 출력
# 여러 개의 슬라이드를 반복하면서 슬라이드를 하나씩 꺼내온다
i=1
for slide in prs.slides:
    # 슬라이드 번호 출력
    print('Slide', i)
    # 하나의 슬라이드로부터 여러개의 쉐이프를 얻은후, 쉐이프를 하나씩 꺼낸다
    for shape in slide.shapes:
        # 쉐이프 타입을 출력한다
        print("Shape: ", shape.shape_type)
    i +=1