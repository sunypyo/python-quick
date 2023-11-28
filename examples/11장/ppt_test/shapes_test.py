import collections
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = '모양 추가'

# 0.93"은 이 전체 모양 세트를 중앙에 배치
left = Inches(0.93)
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

shape = shapes.add_shape(MSO_SHAPE.PENTAGON,left, top, width, height)
shape.text = '1 단계'

left = left + width - Inches(0.4)
# chevrons 모양은 시각적 균형을 위해 더 넓은 너비가 필요하다
width = Inches(2.0)  # chevrons need more width for visual balance

for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = '%d 단계' % n
    left = left + width - Inches(0.4)

prs.save('shapes_test.pptx')