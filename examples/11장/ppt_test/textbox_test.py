import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
tf = txBox.text_frame

tf.text = "안쪽 텍스트 박스 추가"

p = tf.add_paragraph()
p.text = "볼드체로 두번째 단락 추가"
p.font.bold = True

p = tf.add_paragraph()
p.text = "큰 글자체로 세번째 단락 추가"
p.font.size = Pt(40)

prs.save('textBox_test.pptx')
