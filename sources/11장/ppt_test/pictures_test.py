import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches

image = 'pandas.png'
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
#슬라이드 추가
slide = prs.slides.add_slide(blank_slide_layout)

# 이미지를 슬라이드에 추가, 코드에서는 2개의 이미지를 추가.  그림의 좌표와 크기
slide.shapes.add_picture(image, Inches(1), Inches(1))

#아래 코드는 img2를 왼쪽에서 2인치, 상단에서 5인치, 너비가 이미지 너비와 같도록 구성
slide.shapes.add_picture(image, Inches(3), Inches(3.5), Inches(2.5), Inches(4))

prs.save('pictures_test.pptx')
