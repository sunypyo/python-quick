# import collections
# import collections.abc
# 위 두줄의 코드는 python-pptx 라이브러리(패키지)의 버전이 1.0.2 이전 일때 필요
# 명령 프롬프트에서 아래 명령으로 python-pptx 의 지정된 버전 설치
# pip install python-pptx==1.0.2
from pptx import Presentation

prs = Presentation()
two_content_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(two_content_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = '2개의 컨텐츠 슬라이드 추가하기'

tf = body_shape.text_frame
tf.text = '첫번째 라인입니다.'

p = tf.add_paragraph()
p.text = '두번째 라인입니다.'
p.level = 1

p = tf.add_paragraph()
p.text = '세번째 라인입니다.'
p.level = 2

prs.save('placeholders_test.pptx')

