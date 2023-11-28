import collections
import collections.abc
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt

prs = Presentation()

chart_data = ChartData()
chart_data.categories = ['영업부', '마케팅부', '개발부']
chart_data.add_series('2020년도', (47.6, 59.8, 67.8))
chart_data.add_series('2021년도', (77.6, 89.9, 79.9))
chart_data.add_series('2022년도', (80.7, 60.6, 90.3))
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '부서별 데이터 컬럼 차트1'
graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)
chart = graphic_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False
#####

slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '부서별 데이터 컬럼 차트2'
graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

chart = graphic_frame.chart
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

data_labels.font.size = Pt(13)
data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
data_labels.position = XL_LABEL_POSITION.INSIDE_END

#####
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '2022년도 부서별 데이터'
chart_data = ChartData()

chart_data.categories = ['영업부', '마케팅부', '개발부', '운영부', '기타부서']
chart_data.add_series('2022년도', (0.123, 0.199, 0.345, 0.123, 0.456))

x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(6), Inches(5.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
print(type(chart.plots[0]))
# 결과 <class 'pptx.chart.plot.PiePlot'>
data_labels.number_format = '0%'
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

prs.save('ppt_chart.pptx')