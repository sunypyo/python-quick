# import collections
# import collections.abc
# 위 두줄의 코드는 python-pptx 라이브러리(패키지)의 버전이 1.0.2 이전 일때 필요
# 명령 프롬프트에서 아래 명령으로 python-pptx 의 지정된 버전 설치
# pip install python-pptx==1.0.2
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
tf = txBox.text_frame

tf.text = "안쪽 텍스트 박스 추가"

p = tf.add_paragraph()
p.text = "볼드체로 두번째 단락 추가"
p.font.bold = True

p = tf.add_paragraph()
p.text = "큰 글자체로 세번째 단락 추가"
p.font.size = Pt(40)

prs.save('textBox_test.pptx')
